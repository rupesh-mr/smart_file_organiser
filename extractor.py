import os
import json
from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from bs4 import BeautifulSoup
import openpyxl
import pptx


def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext in ['.txt', '.md', '.py']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        elif ext == '.pdf':
            reader = PdfReader(file_path)
            return "\n".join([page.extract_text() or "" for page in reader.pages])

        elif ext == '.docx':
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == '.csv':
            df = pd.read_csv(file_path)
            return df.to_string()

        elif ext == '.html':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f, "html.parser")
                return soup.get_text()

        elif ext == '.json':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)

        elif ext == '.xlsx':
            wb = openpyxl.load_workbook(file_path, read_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            return text

        elif ext == '.pptx':
            prs = pptx.Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            return ""

    except Exception as e:
        print(f"Failed to extract text from {file_path}: {e}")
        return ""
